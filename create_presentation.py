from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Define Constants
IMAGE_DIR = os.path.join(os.getcwd(), 'images')
TITLE_BG = os.path.join(IMAGE_DIR, 'title_background.png')
AR_IMG = os.path.join(IMAGE_DIR, 'ar_crime_scene.png')
VR_IMG = os.path.join(IMAGE_DIR, 'vr_forensics_lab.png')
AI_IMG = os.path.join(IMAGE_DIR, 'ai_data_analysis.png')

DARK_BG_COLOR = RGBColor(10, 15, 30) # Very dark blue
TEXT_COLOR = RGBColor(255, 255, 255) # White
ACCENT_COLOR = RGBColor(0, 191, 255) # Deep Sky Blue

def set_slide_background_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def format_text_frame(text_frame, font_size=Pt(20), bold=False, color=TEXT_COLOR):
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = font_size
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = 'Calibri'
        
def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Set background image
    if os.path.exists(TITLE_BG):
        left = top = 0
        slide.shapes.add_picture(TITLE_BG, left, top, width=prs.slide_width, height=prs.slide_height)
    else:
        set_slide_background_color(slide, DARK_BG_COLOR)

    # Add Title
    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "NexForensics"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Add Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    tf_sub = subtitle.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "AR-Based Forensic Platform & Crime Scene Reconstruction"
    p_sub.font.size = Pt(32)
    p_sub.font.color.rgb = ACCENT_COLOR
    p_sub.alignment = PP_ALIGN.CENTER

def create_content_slide(prs, title_text, content_list, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
    set_slide_background_color(slide, DARK_BG_COLOR)

    # Format Title
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    title.text_frame.paragraphs[0].font.name = 'Calibri Light'
    title.text_frame.paragraphs[0].font.bold = True

    # Format Body
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content_list[0] if content_list else ""
    
    # Formatting first paragraph
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_COLOR

    for item in content_list[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(14)
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
    
    # Add Image if provided
    if image_path and os.path.exists(image_path):
        # Resize body text to make room for image
        body_shape.width = Inches(5.0) 
        
        # Add Image
        left = Inches(5.5)
        top = Inches(2.0)
        height = Inches(4.0)
        slide.shapes.add_picture(image_path, left, top, height=height)

def main():
    prs = Presentation()

    # 1. Title Slide
    add_title_slide(prs)

    # 2. Introduction (with AR Image)
    create_content_slide(prs, "Introduction", [
        "Revolutionizing crime scene investigation.",
        "An AR-based platform for investigators to virtually revisit scenes.",
        "Precision-preserved digital environments.",
        "Merges computer vision with immersive technology."
    ], image_path=AR_IMG)

    # 3. The Problem
    create_content_slide(prs, "The Problem", [
        "Fragile Integrity: Physical crime scenes are easily contaminated.",
        "Juror Confusion: 2D photos lack spatial context and depth.",
        "Slow Analysis: Manual cataloging takes days and is prone to error.",
        "Data Silos: Disconnected tools for evidence gathering and reporting."
    ])

    # 4. The Solution: NexForensics (with VR Image)
    create_content_slide(prs, "The Solution", [
        "AI-Driven Backend: Auto-identification and classification of evidence.",
        "3D Reconstruction: 'Freeze' the crime scene in time with sub-mm accuracy.",
        "Seamless Workflow: From field investigation to classroom or courtroom.",
        "Immersive Tech: Uses AR for field work and VR for analysis."
    ], image_path=VR_IMG)

    # 5. Key Features
    create_content_slide(prs, "Key Features", [
        "Real-Time Overlays: Data augmentation on live scenes.",
        "Remote Consultation: Experts can 'teleport' to the scene via VR.",
        "Chain of Custody: Blockchain-backed evidence tracking.",
        "Automated Discovery: Machine learning highlights potential evidence."
    ])

    # 6. Technology Stack (with AI Image)
    create_content_slide(prs, "Technology Stack", [
        "Frontend: Unity / Unreal Engine for high-fidelity rendering.",
        "Backend: Python, TensorFlow/PyTorch for analysis.",
        "Cloud: AWS/Azure for secure storage and processing.",
        "Hardware: Compatible with HoloLens, Quest 3, and LiDAR tablets."
    ], image_path=AI_IMG)

    # 7. Target Markets
    create_content_slide(prs, "Target Audience", [
        "Law Enforcement Agencies",
        "Forensic Science Universities",
        "Legal Defense & Prosecution Teams",
        "Private Investigators"
    ])

    # 8. Conclusion
    create_content_slide(prs, "Vision", [
        "To become the global standard for digital forensic evidence.",
        "Bridge the gap between physical reality and digital precision.",
        "Empower justice with unquestionable accuracy."
    ])

    output_file = 'NexForensics_Pro_Presentation.pptx'
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

if __name__ == "__main__":
    main()
